#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_pptx — turn the extracted layout manifest into a PowerPoint a designer can work in.

WHAT THIS IS FOR: handing the deck to a designer for layout. So the priority is not
"looks identical in PowerPoint" but "every string is a real, restylable text frame in
the right place, and everything that cannot be text arrives as a clean asset".

    text        → a text box at the measured position, with runs (bold/italic) preserved
    <img>       → the original PNG, embedded at its measured box
    svg/canvas  → the element screenshotted with a transparent background
    KaTeX math  → its TeX SOURCE inside the text (\(x\) / \[x\]) — PowerPoint cannot set
                  KaTeX, and a picture would freeze the surrounding paragraph too; the
                  reference render (ref/sNNN.png) shows how the formula should look
    slide-notes → PowerPoint speaker notes

Positions come from the browser in the deck's design space (1920×1080) and map 1:1 onto a
13.333×7.5 in slide, i.e. 1 design px = 1/144 in. Nothing is re-flowed or guessed.

Usage:
    _research/.venv-pptx/bin/python scripts/build_pptx.py [<deck-stem> ...] [--out DIR]
"""
import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(ROOT, "_internal", "pptx-export")
EMU_PER_PX = 914400 // 144          # 1 design px = 1/144 in — a 1920px slide is 13.333in
ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
         "justify": PP_ALIGN.JUSTIFY}


def px(v):
    return Emu(int(round(v * EMU_PER_PX)))


def add_text(slide, b):
    r, st = b["rect"], b["style"]
    # a little slack: PowerPoint's text metrics are not the browser's, and a box that is
    # exactly as wide as the measured line will wrap one word onto a second line.
    box = slide.shapes.add_textbox(px(r["x"] - 4), px(r["y"] - 3),
                                   px(r["w"] + 14), px(r["h"] + 8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = ALIGN.get(st.get("align"), PP_ALIGN.LEFT)
    if st.get("lh"):
        p.line_spacing = Pt(st["lh"] * 0.75)
    first = True
    for run in b["runs"]:
        t = run["t"]
        for j, chunk in enumerate(t.split("\n")):
            if j:
                p = tf.add_paragraph()
                p.alignment = ALIGN.get(st.get("align"), PP_ALIGN.LEFT)
            if not chunk:
                continue
            rr = p.add_run()
            rr.text = chunk.upper() if st.get("upper") else chunk
            f = rr.font
            f.size = Pt(max(6.0, st["size"] * 0.75))     # css px → pt
            f.name = st.get("family") or "Inter"
            f.bold = bool(run.get("b")) or st.get("weight", 400) >= 600
            f.italic = bool(run.get("i")) or st.get("italic")
            c = st.get("color") or [17, 17, 17]
            f.color.rgb = RGBColor(*[max(0, min(255, int(v))) for v in c])
            first = False
    return box


def add_picture(slide, path, r, note=None):
    if not os.path.exists(path):
        return None
    try:
        return slide.shapes.add_picture(path, px(r["x"]), px(r["y"]),
                                        px(r["w"]), px(r["h"]))
    except Exception as e:                                  # noqa: BLE001
        print(f"    ! картинка не вставлена ({os.path.basename(path)}): {e}")
        return None


def build(stem, outdir):
    d = os.path.join(SRC, stem)
    man_path = os.path.join(d, "manifest.json")
    if not os.path.exists(man_path):
        print(f"  ✗ нет манифеста для {stem} — сначала extract")
        return None
    man = json.load(open(man_path, encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = px(man["design"]["w"])
    prs.slide_height = px(man["design"]["h"])
    blank = prs.slide_layouts[6]                            # blank: no placeholder junk

    n_text = n_img = n_fig = n_missing = 0
    for s in man["slides"]:
        sl = prs.slides.add_slide(blank)
        tex_notes = []
        for b in s["blocks"]:
            if b["kind"] == "text":
                add_text(sl, b)
                n_text += 1
            elif b["kind"] == "image":
                p = os.path.join(ROOT, "Lectures", b["src"])
                if not os.path.exists(p):                   # build output ships .webp
                    alt = os.path.splitext(p)[0] + ".webp"
                    p = alt if os.path.exists(alt) else p
                if add_picture(sl, p, b["rect"]) is None:
                    n_missing += 1
                else:
                    n_img += 1
            elif b["kind"] == "figure":
                if b.get("file") and add_picture(sl, os.path.join(d, b["file"]), b["rect"]):
                    n_fig += 1
                else:
                    n_missing += 1
                if b.get("tex"):
                    tex_notes.append(b["tex"])
        # speaker notes: the deck's own notes + the source of any math on the slide
        parts = [f'[{s["label"]}] тип: {s["type"]}']
        if s.get("notes"):
            parts.append(s["notes"])
        if tex_notes:
            parts.append("Формулы/фигуры (исходный текст): " + " | ".join(tex_notes))
        sl.notes_slide.notes_text_frame.text = "\n\n".join(parts)

    os.makedirs(outdir, exist_ok=True)
    out = os.path.join(outdir, f"{stem}.pptx")
    prs.save(out)
    print(f"  ✓ {stem}.pptx — {len(man['slides'])} слайдов · {n_text} текстовых блоков · "
          f"{n_img} иллюстраций · {n_fig} фигур"
          + (f" · {n_missing} не вставлено" if n_missing else ""))
    return out


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("stems", nargs="*", help="например 13-crucible-of-negatives (по умолчанию — все извлечённые)")
    ap.add_argument("--out", default=os.path.join(ROOT, "_internal", "pptx-export", "_pptx"))
    a = ap.parse_args()
    stems = a.stems or sorted(x for x in os.listdir(SRC)
                              if os.path.isdir(os.path.join(SRC, x)) and not x.startswith("_"))
    if not stems:
        print("нечего собирать: сначала node scripts/export-pptx-extract.mjs")
        return 1
    print(f"[build_pptx] {len(stems)} дек(ов) → {a.out}")
    ok = 0
    for s in stems:
        if build(s, a.out):
            ok += 1
    print(f"[build_pptx] готово: {ok}/{len(stems)}")
    return 0 if ok == len(stems) else 1


if __name__ == "__main__":
    sys.exit(main())
